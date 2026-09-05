#!/usr/bin/env python3
"""
Office Document Reader - docx, xlsx, pptx, doc, xls, ppt
"""

import argparse
import os
import subprocess
import re
from pathlib import Path

WORK_TMP_DIR = './work/tmp'

def ensure_tmp_dir():
    os.makedirs(WORK_TMP_DIR, exist_ok=True)

def convert_with_libreoffice(input_path):
    ext = Path(input_path).suffix.lower()
    convert_map = {
        '.doc': 'docx',
        '.xls': 'xlsx',
        '.ppt': 'pptx',
    }
    target_ext = convert_map.get(ext)
    if not target_ext:
        return input_path

    ensure_tmp_dir()
    cmd = [
        'soffice', '--headless', '--convert-to', target_ext,
        '--outdir', WORK_TMP_DIR, input_path
    ]
    try:
        subprocess.run(cmd, check=True, capture_output=True, timeout=60)
    except Exception:
        print(f'[convert] failed')
        return input_path

    base_name = Path(input_path).stem
    converted = os.path.join(WORK_TMP_DIR, f'{base_name}.{target_ext}')
    if os.path.exists(converted):
        print(f'[convert] ok: {converted}')
        return converted
    return input_path


def read_docx(path):
    from docx import Document
    doc = Document(path)

    # 正文
    body = '\n'.join(p.text for p in doc.paragraphs).rstrip('\n')

    # 表格
    tables = ''
    for table in doc.tables:
        for row in table.rows:
            seen = set()
            cells = []
            for cell in row.cells:
                # 合并单元格延续位置留空，保持列对齐
                if id(cell._tc) in seen:
                    cells.append('')
                else:
                    seen.add(id(cell._tc))
                    cells.append(cell.text.strip().replace('\n', '\\n'))
            row_text = ' | '.join(cells)
            if row_text:
                tables += row_text + '\n'
        tables += '\n'

    # 批注
    comments = ''
    for c in getattr(doc, 'comments', []):
        comments += c.text + '\n'

    return {'body': body, 'tables': tables, 'comments': comments}


def read_xlsx(path):
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True)
    body_parts, comments = [], []
    for sheet in wb:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            if cells := [str(c) if c else '' for c in row]:
                rows.append(' | '.join(cells))
        if rows:
            body_parts.append(f"[Sheet: {sheet.title}]\n" + '\n'.join(rows))
        for row in sheet.iter_rows():
            for cell in row:
                if cell.comment:
                    comments.append(f"{cell.comment.author}{cell.comment.text}")
    body = '\n\n'.join(body_parts)
    return {'body': body, 'tables': body, 'comments': '\n'.join(comments)}


def read_pptx(path):
    from pptx import Presentation
    import zipfile

    prs = Presentation(path)
    slides, comments = [], []
    tables_list = []

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for s in slide.shapes:
            if hasattr(s, 'text') and s.text.strip():
                texts.append(s.text)
            if hasattr(s, 'has_table') and s.has_table:
                rows = []
                for row in s.table.rows:
                    seen = set()
                    cells = []
                    for cell in row.cells:
                        # 合并单元格延续位置留空，保持列对齐
                        if id(cell._tc) in seen:
                            cells.append('')
                        else:
                            seen.add(id(cell._tc))
                            cells.append(cell.text.strip().replace('\n', '\\n'))
                    if cells:
                        rows.append(' | '.join(cells))
                tables_list.append(f"[Slide {i}] " + '\n'.join(rows))
        if texts:
            slides.append(f"[Slide {i}] " + ' '.join(texts))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                comments.append(f"[Slide {i}] {notes}")

    try:
        with zipfile.ZipFile(path, 'r') as z:
            for name in z.namelist():
                if 'comments/modernComment' in name:
                    content = z.read(name).decode('utf-8')
                    texts = re.findall(r'<a:t>([^<]+)</a:t>', content)
                    if texts:
                        comments.append(' '.join(texts))
    except Exception:
        pass

    return {'body': '\n'.join(slides), 'tables': '\n\n'.join(tables_list), 'comments': '\n'.join(comments)}


def read_file(path):
    converted_path = convert_with_libreoffice(path)
    ext = Path(converted_path).suffix.lower()

    readers = {
        '.docx': read_docx,
        '.xlsx': read_xlsx,
        '.pptx': read_pptx,
    }
    reader = readers.get(ext)
    if not reader:
        return {'error': f'Unsupported: {ext}'}

    result = reader(converted_path)
    if converted_path != path:
        result['converted_from'] = path
    return result


def print_result(result, mode):
    if mode in result and result[mode]:
        print(result[mode])


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('file')
    ap.add_argument('--body', action='store_true')
    ap.add_argument('--comments', action='store_true')
    ap.add_argument('--tables', action='store_true')
    args = ap.parse_args()

    result = read_file(args.file)
    if 'error' in result:
        print(result['error'])
        exit(1)

    modes = []
    if args.body:
        modes.append('body')
    if args.comments:
        modes.append('comments')
    if args.tables:
        modes.append('tables')

    if not modes:
        modes = ['body', 'comments', 'tables']

    for mode in modes:
        print(f'=== {mode.capitalize()} ===')
        print_result(result, mode)
        print()


if __name__ == '__main__':
    main()