#!/usr/bin/env python3
"""
Office Document Modifier - docx, xlsx, pptx, doc, xls, ppt

功能说明:
    - 使用 LibreOffice/soffice 将 legacy Office 格式 (doc/xls/ppt) 转换为现代格式 (docx/xlsx/pptx)
    - 读取现代 Office 格式 (docx/xlsx/pptx) 使用 python-docx/openpyxl/pptx 库
    - 修改正文内容、表格中的文本
"""

import argparse
import os
import subprocess
import shutil
import zipfile
import tempfile
from pathlib import Path

WORK_TMP_DIR = './work/tmp'

def ensure_tmp_dir():
    os.makedirs(WORK_TMP_DIR, exist_ok=True)

def convert_with_libreoffice(input_path):
    ext = Path(input_path).suffix.lower()
    convert_map = {
        '.doc': 'docx',
        '.xls': 'xlsx',
        '.ppt': 'pptx',
    }
    target_ext = convert_map.get(ext)
    if not target_ext:
        return input_path

    ensure_tmp_dir()
    cmd = [
        'soffice', '--headless', '--convert-to', target_ext,
        '--outdir', WORK_TMP_DIR, input_path
    ]
    try:
        subprocess.run(cmd, check=True, capture_output=True, timeout=60)
    except Exception:
        return input_path

    base_name = Path(input_path).stem
    converted = os.path.join(WORK_TMP_DIR, f'{base_name}.{target_ext}')
    return converted if os.path.exists(converted) else input_path


def modify_docx(path, old_text, new_text, output_path):
    from docx import Document

    doc = Document(path)
    modified = False

    for para in doc.paragraphs:
        if old_text in para.text:
            para.text = para.text.replace(old_text, new_text)
            modified = True

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if old_text in cell.text:
                    cell.text = cell.text.replace(old_text, new_text)
                    modified = True

    for shape in doc.inline_shapes:
        if hasattr(shape, 'text') and old_text in shape.text:
            shape.text = shape.text.replace(old_text, new_text)
            modified = True

    doc.save(output_path)
    return modified


def modify_xlsx(path, old_text, new_text, output_path):
    from openpyxl import load_workbook

    wb = load_workbook(path)
    modified = False

    for sheet in wb:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str) and old_text in cell.value:
                    cell.value = cell.value.replace(old_text, new_text)
                    modified = True

    wb.save(output_path)
    return modified


def modify_pptx(path, old_text, new_text, output_path):
    from pptx import Presentation

    prs = Presentation(path)
    modified = False

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and old_text in shape.text:
                shape.text = shape.text.replace(old_text, new_text)
                modified = True

    prs.save(output_path)
    return modified


def modify_file(input_path, old_text, new_text, output_path):
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    converted_path = convert_with_libreoffice(input_path)
    ext = Path(converted_path).suffix.lower()

    modifiers = {
        '.docx': modify_docx,
        '.xlsx': modify_xlsx,
        '.pptx': modify_pptx,
    }

    modifier = modifiers.get(ext)
    if not modifier:
        return False

    # modifier 已直接把修改结果保存到 output_path，
    # 无需再复制（旧代码里的 shutil.copy(self) 是无效自复制，已移除）
    result = modifier(converted_path, old_text, new_text, output_path)

    return result


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('file', help='Input file path')
    ap.add_argument('--old', required=True, help='Text to replace')
    ap.add_argument('--new', required=True, help='New text')
    ap.add_argument('--output', required=True, help='Output file path')
    args = ap.parse_args()

    success = modify_file(args.file, args.old, args.new, args.output)

    if success:
        print(f'Modified: {args.output}')
    else:
        print('No modification made')


if __name__ == '__main__':
    main()